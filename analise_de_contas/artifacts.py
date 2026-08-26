"""
Módulo de geração de artefatos — PPTX fallback, HTML animado, CSV, consolidação.
"""

import os
import re
import csv
import json
import logging
from pathlib import Path

from jinja2 import Environment, FileSystemLoader, select_autoescape

_logger = logging.getLogger("pipeline")


def log(msg: str):
    _logger.info(msg)


# ===== CSV =====

def convert_to_csv(md_file: str, csv_file: str) -> int:
    """Converte tabela Markdown para CSV. Retorna número de linhas escritas.
    
    Lida com células multilinha (continuação de tabela sem | no início).
    """
    with open(md_file, "r", encoding="utf-8") as f:
        content = f.read()

    lines = content.split("\n")
    table_lines = []
    in_table = False
    current_row = ""

    for line in lines:
        stripped = line.strip()
        # Nova linha de tabela começando com |
        if stripped.startswith("|") and "|" in stripped[1:]:
            if re.match(r"\|[\s\-:|]+\|", stripped):
                continue  # separator row
            if in_table:
                if current_row:
                    table_lines.append(current_row)
                current_row = stripped
            else:
                in_table = True
                current_row = stripped
        elif in_table and stripped:
            # Continuação de célula multilinha — concatena
            if current_row:
                # Anexa ao último célula da linha atual
                current_row = current_row.rstrip("|") + " " + stripped + "|"
        elif in_table:
            if current_row:
                table_lines.append(current_row)
            current_row = ""
            in_table = False

    if in_table and current_row:
        table_lines.append(current_row)

    if not table_lines:
        return 0

    rows = []
    for line in table_lines:
        cells = [c.strip() for c in line.split("|")[1:-1]]
        if len(cells) >= 5:
            rows.append(cells)

    if rows:
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in rows:
                writer.writerow(row)
        return len(rows)
    return 0


# ===== Markdown Section Extraction =====

def extract_md_sections(md_file: str) -> dict:
    """Extrai seções de um arquivo Markdown para uso no fallback PPTX."""
    if not os.path.exists(md_file):
        return {}
    with open(md_file, "r", encoding="utf-8") as f:
        content = f.read()

    sections = {}
    current_key = "intro"
    current_lines = []

    for line in content.split("\n"):
        if line.startswith("# "):
            if current_lines:
                sections[current_key] = "\n".join(current_lines).strip()
            current_key = line[2:].strip().lower().replace(" ", "_")[:30]
            current_lines = [line]
        else:
            current_lines.append(line)

    if current_lines:
        sections[current_key] = "\n".join(current_lines).strip()

    return sections


def _read_section(output_dir: Path, path: str, limit: int = 800) -> str:
    """Lê uma seção de relatório Markdown com limite de caracteres."""
    p = output_dir / path
    return p.read_text(encoding="utf-8")[:limit] if p.exists() else ""


def _split_bullets(text: str, max_items: int = 6) -> list[str]:
    """Divide texto em bullets, cortando em sentenças."""
    return [s.strip() + "." for s in text.split(". ") if s.strip()][:max_items]


def _strip_nb_preamble(text: str) -> str:
    """Remove o preamble do NotebookLM (primeiras 3 linhas boilerplate)."""
    lines = text.split("\n")
    start = 0
    for i, line in enumerate(lines):
        if line.startswith("# ") or line.startswith("| ") or line.startswith("- "):
            start = i
            break
    return "\n".join(lines[start:])


# ===== PPTX Fallback =====

def generate_pptx_fallback(output_dir: Path, pdf_name: str, slide_type: str) -> bool:
    """Gera apresentação PPTX localmente via python-pptx como fallback."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        log("  ✗ python-pptx não instalado — fallback PPTX indisponível")
        return False

    NAVY = RGBColor(0x1A, 0x23, 0x7E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    RED = RGBColor(0xC6, 0x28, 0x28)
    ORANGE = RGBColor(0xE6, 0x51, 0x00)
    LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
    DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
    ACCENT = RGBColor(0x28, 0x35, 0x93)
    MUTED = RGBColor(0x79, 0x86, 0xCB)

    parecer = _read_section(output_dir, "06_parecer_final.md")
    metodologia = _read_section(output_dir, "01_receitas.md")
    gaps = _read_section(output_dir, "04_consistencia_financeira.md")
    escrita = _read_section(output_dir, "05_qualidade_documental.md", 1500)
    if "RESUMO ESTATÍSTICO" in escrita:
        erros = escrita.split("RESUMO ESTATÍSTICO")[-1][:500]
    else:
        erros = escrita[-500:]

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title(slide, text, color=NAVY, y=0.3):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.color.rgb = color
        p.font.bold = True

    def add_bullets(slide, items, x=0.5, y=1.3, w=9, h=3.8, color=DARK_TEXT):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = color
            p.level = 0
            p.space_after = Pt(6)

    def add_accent_line(slide, y=1.0, color=ACCENT):
        shape = slide.shapes.add_shape(
            1, Inches(0.5), Inches(y), Inches(9), Inches(0.05)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]

    if slide_type == "completa":
        prs.core_properties.title = "Auditoria de Contas Condominiais"
        prs.core_properties.author = "AnaliseContas v7.0 + NotebookLM"

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, NAVY)
        txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "AUDITORIA DE CONTAS CONDOMINAIS"
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        txBox2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.8))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = pdf_name
        p2.font.size = Pt(20)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER
        txBox3 = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
        p3 = txBox3.text_frame.paragraphs[0]
        p3.text = "Gerado via AnaliseContas v7.0. Atenção: Esta é uma análise gerada por IA e pode estar sujeita a erros."
        p3.font.size = Pt(12)
        p3.font.color.rgb = MUTED
        p3.alignment = PP_ALIGN.CENTER

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "RESUMO EXECUTIVO")
        add_bullets(s, _split_bullets(parecer))

        s = prs.slides.add_slide(blank_layout)
        add_title(s, "RECEITAS")
        add_accent_line(s)
        add_bullets(s, _split_bullets(metodologia))

        s = prs.slides.add_slide(blank_layout)
        add_title(s, "INCONSISTÊNCIAS FINANCEIRAS", RED)
        add_accent_line(s, color=RED)
        add_bullets(s, _split_bullets(gaps))

        s = prs.slides.add_slide(blank_layout)
        add_title(s, "QUALIDADE DOCUMENTAL")
        add_bullets(s, _split_bullets(erros))

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, NAVY)
        txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "PARECER FINAL"
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        txBox2 = s.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = "AnaliseContas v7.0. Atenção: Esta é uma análise gerada por IA e pode estar sujeita a erros."
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER

        out_file = output_dir / "apresentacao_completa.pptx"

    else:  # auditoria
        prs.core_properties.title = "Relatório de Auditoria Condominial"
        prs.core_properties.author = "AnaliseContas v7.0 + NotebookLM"

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, RED)
        txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "RELATÓRIO DE AUDITORIA CONDOMINIAL"
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        txBox2 = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.8))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = pdf_name
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
        p2.alignment = PP_ALIGN.CENTER
        
        txBox3 = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
        p3 = txBox3.text_frame.paragraphs[0]
        p3.text = "Atenção: Esta é uma análise gerada por IA e pode estar sujeita a erros."
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
        p3.alignment = PP_ALIGN.CENTER

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "INCONSISTÊNCIAS FINANCEIRAS", RED)
        add_accent_line(s, color=RED)
        add_bullets(s, _split_bullets(gaps))

        s = prs.slides.add_slide(blank_layout)
        add_title(s, "PROBLEMAS DOCUMENTAIS", ORANGE)
        add_accent_line(s, color=ORANGE)
        add_bullets(s, _split_bullets(erros))

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, NAVY)
        txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "RECOMENDAÇÕES"
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        add_bullets(s, _split_bullets(parecer), y=2.2, h=3, color=RGBColor(0xE8, 0xEA, 0xF6))

        out_file = output_dir / "apresentacao_auditoria.pptx"

    try:
        prs.save(str(out_file))
        size = out_file.stat().st_size
        if size > 1000:
            log(f"  ✓ PPTX fallback gerado: {size} bytes (python-pptx)")
            return True
        else:
            log(f"  ✗ PPTX muito pequeno: {size} bytes")
            return False
    except Exception as e:
        log(f"  ✗ Erro ao salvar PPTX: {e}")
        return False


# ===== HTML Animado (Mira) =====

def generate_mira_artifact(output_dir: Path, pdf_name: str) -> bool:
    """Gera apresentação animada HTML via padrões mira-animator (D3.js + Tailwind)."""

    parecer_file = output_dir / "06_parecer_final.md"
    receitas_file = output_dir / "01_receitas.md"
    consistencia_file = output_dir / "04_consistencia_financeira.md"
    estrutura_file = output_dir / "00_estrutura_documento.md"
    legal_file = output_dir / "02_conformidade_legal.md"
    despesas_file = output_dir / "03_despesas.md"
    qualidade_file = output_dir / "05_qualidade_documental.md"

    parecer = parecer_file.read_text(encoding="utf-8") if parecer_file.exists() else ""
    metodologia = receitas_file.read_text(encoding="utf-8") if receitas_file.exists() else ""
    gaps = consistencia_file.read_text(encoding="utf-8") if consistencia_file.exists() else ""
    estrutura = estrutura_file.read_text(encoding="utf-8") if estrutura_file.exists() else ""
    auditoria = legal_file.read_text(encoding="utf-8") if legal_file.exists() else ""
    escrita = qualidade_file.read_text(encoding="utf-8") if qualidade_file.exists() else ""

    from utils import extract_structured_score
    structured = extract_structured_score(parecer)

    nota = str(structured["nota"]) if structured.get("nota") is not None else "N/A"
    decisao = structured.get("decisao") or "N/A"
    _nota_found = nota != "N/A"
    _dec_found = decisao != "N/A"

    for line in parecer.split("\n"):
        # NOTA: regex permissivo — captura qualquer variação de formatação
        if not _nota_found:
            m_nota = re.search(r'(?i)\bNOTA\s*(?:FINAL)?\b[^0-9]*(\d+[.,]?\d*)', line)
            if m_nota and m_nota.group(1):
                nota = m_nota.group(1)
                _nota_found = True
        # DECISÃO: regex permissivo (Accept, Minor/Major Revisions, Reject, etc.)
        if not _dec_found:
            m_dec = re.search(r'(?i)\bDECIS[ÃA]O\s*(?:EDITORIAL|FINAL|RECOMENDADA)?\s*[:=\-—]?\s*(?:\*\*\s*)?(.+?)(?:\s*\*\*)?$', line)
            if m_dec and m_dec.group(1).strip():
                decisao = re.sub(r'\*\*', '', m_dec.group(1)).strip()
                _dec_found = True

    gaps_list = []
    for line in gaps.split("\n"):
        match = re.match(r"^(?:- |\d+\.\s+)(.+)", line.strip())
        if match:
            gap_text = match.group(1).replace("**", "")
            if ":" in gap_text:
                gap_text = gap_text.split(":", 1)[0]
            gaps_list.append(gap_text.strip()[:80])
    gaps_list = gaps_list[:5]

    fortes_list = list(structured.get("pontos_fortes", []))[:4]
    if not fortes_list:
        in_fortes = False
        for line in parecer.split("\n"):
            if re.match(r'(?i)^#{1,3}\s*(?:\d+\.?\d*\s*)?(?:\*\*)?(?:pontos?\s+fortes?|for[çc]as?\s+(?:do|identificad))', line):
                in_fortes = True
                continue
            if in_fortes and re.match(r'^#{1,3}\s', line):
                in_fortes = False

            if in_fortes:
                match = re.match(r"^(?:- |\d+\.\s+)(.+)", line.strip())
                if match:
                    forte_text = match.group(1).replace("**", "")
                    if ":" in forte_text:
                        forte_text = forte_text.split(":", 1)[0]
                    fortes_list.append(forte_text.strip()[:80])
        fortes_list = fortes_list[:4]

    # --- Extract Fragilidades from parecer ---
    fragilidades_list = list(structured.get("fragilidades", []))[:4]
    if not fragilidades_list:
        in_fragilidades = False
        for line in parecer.split("\n"):
            if re.match(r'(?i)^#{1,3}\s*(?:\d+\.?\d*\s*)?(?:\*\*)?(?:fragilidades?|pontos?\s+fracos?|limita[çc][õo]es?|fraquezas?|weakness)', line):
                in_fragilidades = True
                continue
            if in_fragilidades and re.match(r'^#{1,3}\s', line):
                in_fragilidades = False
            if in_fragilidades:
                match = re.match(r"^(?:- |\d+\.\s+)(.+)", line.strip())
                if match:
                    texto = match.group(1).replace("**", "")
                    if ":" in texto:
                        texto = texto.split(":", 1)[0]
                    fragilidades_list.append(texto.strip()[:80])
        fragilidades_list = fragilidades_list[:4]

    # --- Extract Resumo Executivo from parecer ---
    resumo_executivo = ""
    in_resumo = False
    resumo_lines = []
    for line in parecer.split("\n"):
        if re.search(r'(?i)resumo\s*executivo', line):
            in_resumo = True
            continue
        if in_resumo and re.match(r'^#{1,3}\s', line):
            in_resumo = False
        if in_resumo and line.strip():
            resumo_lines.append(line.strip())
    resumo_executivo = " ".join(resumo_lines)[:500]

    # --- Extract methodology steps dynamically from 01_receitas.md ---
    metodo_steps = []
    for line in metodologia.split("\n"):
        m = re.match(r'^- \*\*(.+?)\*\*', line.strip())
        if m:
            metodo_steps.append(m.group(1).strip()[:50])
    if not metodo_steps:
        for line in metodologia.split("\n"):
            m = re.match(r'^(?:\d+[.)]\s+)(.+)', line.strip())
            if m:
                metodo_steps.append(m.group(1).replace("**", "").strip()[:50])
    if not metodo_steps:
        for line in metodologia.split("\n"):
            m = re.match(r'^#{2,4}\s+(.+)', line.strip())
            if m:
                metodo_steps.append(m.group(1).replace("**", "").strip()[:50])
    if not metodo_steps:
        for line in metodologia.split("\n"):
            m = re.match(r'^(?:[-•*>])\s+(.+)', line.strip())
            if m:
                text = m.group(1).replace("**", "").strip()[:50]
                if len(text) > 5:
                    metodo_steps.append(text)
    metodo_steps = metodo_steps[:5] or ["Taxas Condominiais", "Fundos", "Multas", "Outras Receitas"]

    # --- Extract writing error count from 05_qualidade_documental.md ---
    escrita_total_erros = 0
    if escrita:
        # Count table rows (lines starting with |)
        table_rows = [l for l in escrita.split("\n") if l.strip().startswith("|") and "Seção" not in l and "---" not in l]
        escrita_total_erros = len(table_rows)
        # Also try to find explicit count
        for line in escrita.split("\n"):
            m = re.search(r'(\d+)\s+erros?', line, re.IGNORECASE)
            if m:
                escrita_total_erros = int(m.group(1))
                break

    # --- Extract audit checklist items from 02_conformidade_legal.md ---
    auditoria_items = []
    for line in auditoria.split("\n"):
        # Formato checkbox: - [x] **Item** ou - [ ] **Item**
        m = re.match(r'^- \[([ xX✓✗])\]\s*\*\*(.+?)\*\*', line.strip())
        if m:
            status = "ok" if m.group(1).lower() in ('x', '✓') else "fail"
            auditoria_items.append({"item": m.group(2).strip()[:40], "status": status})
            continue
        # Formato emoji: - ✅ **Item** ou - ❌ **Item**
        m2 = re.match(r'^[-•]\s*([✅✓☑❌✗☒⚠️])\s*\*?\*?(.+?)(?:\*\*)?$', line.strip())
        if m2:
            status = "ok" if m2.group(1) in ('✅', '✓', '☑') else "fail"
            auditoria_items.append({"item": m2.group(2).replace("**", "").strip()[:40], "status": status})
            continue
        # Formato tabela: | Item | Sim/Não |
        m3 = re.match(r'^\|(.+?)\|(.+?)\|', line.strip())
        if m3 and '---' not in m3.group(1):
            item_text = m3.group(1).replace("**", "").strip()
            status_text = m3.group(2).strip().lower()
            if item_text and len(item_text) > 2:
                status = "ok" if any(s in status_text for s in ('sim', 'yes', '✅', '✓', 'ok', 'presente')) else "fail"
                auditoria_items.append({"item": item_text[:40], "status": status})
    auditoria_items = auditoria_items[:6]

    # Diagnóstico: log dos dados extraídos para a apresentação
    log(f"  📊 Dados extraídos para apresentação HTML:")
    log(f"    NOTA: {nota} | DECISÃO: {decisao}")
    log(f"    Gaps: {len(gaps_list)} | Fortes: {len(fortes_list)} | Fragilidades: {len(fragilidades_list)}")
    log(f"    Metodologia steps: {len(metodo_steps)} | Auditoria items: {len(auditoria_items)}")
    log(f"    Resumo executivo: {len(resumo_executivo)} chars | Erros escrita: {escrita_total_erros}")
    if nota == "N/A":
        log("    ⚠️ NOTA não extraída — verifique formato de 06_parecer_final.md")
    if decisao == "N/A":
        log("    ⚠️ DECISÃO não extraída — verifique formato de 06_parecer_final.md")

    slug = re.sub(r"[^a-z0-9]", "", pdf_name.lower())[:20] or "paper"
    gaps_json = json.dumps(gaps_list[:5] if gaps_list else ["Sem dados"])
    nota_clean = nota.replace(",", ".")
    decisao_short = decisao[:30]
    decisao_60 = decisao[:60]

    # Carregar template Jinja2
    template_dir = Path(__file__).resolve().parent / "templates"
    env = Environment(loader=FileSystemLoader(str(template_dir)), autoescape=select_autoescape(["html", "xml"]))
    template = env.get_template("mira.html")

    html = template.render(
        slug=slug,
        pdf_name=pdf_name,
        nota=nota,
        nota_clean=nota_clean,
        decisao_short=decisao_short,
        decisao_60=decisao_60,
        gaps_json=gaps_json,
        gaps_list=gaps_list,
        fortes_json=json.dumps(fortes_list if fortes_list else ["Nenhum ponto forte destacado"]),
        fortes_list=fortes_list,
        fragilidades_list=fragilidades_list,
        fragilidades_json=json.dumps(fragilidades_list if fragilidades_list else []),
        metodo_steps_json=json.dumps(metodo_steps),
        escrita_total_erros=escrita_total_erros,
        escrita_text=_strip_nb_preamble(escrita)[:400] if escrita else "Qualidade documental não disponível.",
        auditoria_items=auditoria_items,
        auditoria_json=json.dumps(auditoria_items if auditoria_items else []),
        resumo_executivo=resumo_executivo if resumo_executivo else "Resumo executivo não disponível.",
        metodologia_text=_strip_nb_preamble(metodologia)[:500] if metodologia else "Receitas não analisadas.",
        total_slides=7,
    )

    output_file = output_dir / "apresentacao_animada.html"
    output_file.write_text(html, encoding="utf-8")
    size = output_file.stat().st_size

    if size > 1000:
        log(f"  ✓ Apresentação animada HTML gerada: {size} bytes")
        return True
    else:
        log(f"  ✗ HTML muito pequeno: {size} bytes")
        return False



