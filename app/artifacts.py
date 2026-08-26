"""
Módulo de geração de artefatos — PPTX fallback, HTML animado, CSV, consolidação.
"""

import os
import re
import csv
import json
import logging
from pathlib import Path

from jinja2 import Environment, FileSystemLoader, select_autoescape

_logger = logging.getLogger("pipeline")


def log(msg: str):
    _logger.info(msg)


# ===== CSV =====

def convert_to_csv(md_file: str, csv_file: str) -> int:
    """Converte tabela Markdown para CSV. Retorna número de linhas escritas.
    
    Lida com células multilinha (continuação de tabela sem | no início).
    """
    with open(md_file, "r", encoding="utf-8") as f:
        content = f.read()

    lines = content.split("\n")
    table_lines = []
    in_table = False
    current_row = ""

    for line in lines:
        stripped = line.strip()
        # Nova linha de tabela começando com |
        if stripped.startswith("|") and "|" in stripped[1:]:
            if re.match(r"\|[\s\-:|]+\|", stripped):
                continue  # separator row
            if in_table:
                if current_row:
                    table_lines.append(current_row)
                current_row = stripped
            else:
                in_table = True
                current_row = stripped
        elif in_table and stripped:
            # Continuação de célula multilinha — concatena
            if current_row:
                # Anexa ao último célula da linha atual
                current_row = current_row.rstrip("|") + " " + stripped + "|"
        elif in_table:
            if current_row:
                table_lines.append(current_row)
            current_row = ""
            in_table = False

    if in_table and current_row:
        table_lines.append(current_row)

    if not table_lines:
        return 0

    rows = []
    for line in table_lines:
        cells = [c.strip() for c in line.split("|")[1:-1]]
        if len(cells) >= 5:
            rows.append(cells)

    if rows:
        with open(csv_file, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in rows:
                writer.writerow(row)
        return len(rows)
    return 0


# ===== Markdown Section Extraction =====

def extract_md_sections(md_file: str) -> dict:
    """Extrai seções de um arquivo Markdown para uso no fallback PPTX."""
    if not os.path.exists(md_file):
        return {}
    with open(md_file, "r", encoding="utf-8") as f:
        content = f.read()

    sections = {}
    current_key = "intro"
    current_lines = []

    for line in content.split("\n"):
        if line.startswith("# "):
            if current_lines:
                sections[current_key] = "\n".join(current_lines).strip()
            current_key = line[2:].strip().lower().replace(" ", "_")[:30]
            current_lines = [line]
        else:
            current_lines.append(line)

    if current_lines:
        sections[current_key] = "\n".join(current_lines).strip()

    return sections


def _read_section(output_dir: Path, path: str, limit: int = 800) -> str:
    """Lê uma seção de relatório Markdown com limite de caracteres."""
    p = output_dir / path
    return p.read_text(encoding="utf-8")[:limit] if p.exists() else ""


def _split_bullets(text: str, max_items: int = 6) -> list[str]:
    """Divide texto em bullets, cortando em sentenças."""
    return [s.strip() + "." for s in text.split(". ") if s.strip()][:max_items]


def _strip_nb_preamble(text: str) -> str:
    """Remove o preamble do NotebookLM (primeiras 3 linhas boilerplate)."""
    lines = text.split("\n")
    start = 0
    for i, line in enumerate(lines):
        if line.startswith("# ") or line.startswith("| ") or line.startswith("- "):
            start = i
            break
    return "\n".join(lines[start:])


# ===== PPTX Fallback =====

def generate_pptx_fallback(output_dir: Path, pdf_name: str, slide_type: str) -> bool:
    """Gera apresentação PPTX completa e profissional localmente via python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        log("  ✗ python-pptx não instalado — fallback PPTX indisponível")
        return False

    NAVY = RGBColor(0x0F, 0x17, 0x2A)       # Slate 900
    DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)  # Slate 800
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    RED = RGBColor(0xEF, 0x44, 0x44)        # Rose 500
    ORANGE = RGBColor(0xF5, 0x9E, 0x0B)     # Amber 500
    GREEN = RGBColor(0x10, 0xB9, 0x81)      # Emerald 500
    LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)   # Slate 50
    DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)  # Slate 800
    TEAL = RGBColor(0x0D, 0x94, 0x88)       # Teal 600
    MUTED = RGBColor(0x64, 0x74, 0x8B)      # Slate 500
    LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1) # Teal 100

    from utils import extract_structured_score
    parecer_raw = _read_section(output_dir, "06_sintese_parecer.md", 5000)
    structured = extract_structured_score(parecer_raw) if parecer_raw else {}

    nota = structured.get("nota")
    nota_str = f"{nota:.1f}/10" if nota is not None else "N/A"
    decisao = structured.get("decisao") or "Pendente"

    # Seções
    estrutura = _read_section(output_dir, "00_estrutura_documento.md", 2000)
    metodologia = _read_section(output_dir, "01_metodologia.md", 2000)
    editorial = _read_section(output_dir, "02_auditoria_editorial.md", 2000)
    sota = _read_section(output_dir, "03_sota_referencias.md", 2000)
    gaps = _read_section(output_dir, "04_gaps_logicos.md", 2000)
    escrita = _read_section(output_dir, "05_analise_escrita.md", 2000)
    quantitativa = _read_section(output_dir, "07_auditoria_quantitativa.md", 2000)
    bibliografia = _read_section(output_dir, "08_auditoria_bibliografica.md", 2000)

    # Erros de escrita
    if "RESUMO ESTATÍSTICO" in escrita:
        erros_text = escrita.split("RESUMO ESTATÍSTICO")[-1][:800]
    else:
        erros_text = escrita[:800]

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title(slide, text, subtitle=None, color=NAVY, y=0.4):
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = color
        p.font.bold = True
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(12)
            p2.font.color.rgb = MUTED
            p2.space_before = Pt(3)

    def add_bullets(slide, items, x=0.6, y=1.4, w=8.8, h=3.8, color=DARK_TEXT, font_size=13):
        if not items:
            items = ["Nenhum dado específico registrado nesta seção."]
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}" if not item.startswith("•") else item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.space_after = Pt(8)

    def add_accent_line(slide, y=1.2, color=TEAL):
        shape = slide.shapes.add_shape(
            1, Inches(0.6), Inches(y), Inches(8.8), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def add_score_badge(slide, x=6.8, y=0.35, w=2.6, h=0.8):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_CARD
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"NOTA: {nota_str} | {decisao}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]

    if slide_type == "completa":
        prs.core_properties.title = "Parecer Técnico & Auditoria Peer-Review"
        prs.core_properties.author = "Hermes Agent + NotebookLM v6.0"

        # 1. Capa
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, NAVY)
        txBox = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(1.2))
        p = txBox.text_frame.paragraphs[0]
        p.text = "PARECER TÉCNICO & AUDITORIA PEER-REVIEW"
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(0.8))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = f"Avaliação Crítica do Paper: {pdf_name}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = LIGHT_TEAL
        p2.alignment = PP_ALIGN.CENTER

        txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.5))
        p3 = txBox3.text_frame.paragraphs[0]
        p3.text = f"Banca Examinadora Acadêmica · Nota: {nota_str} · Decisão: {decisao}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = MUTED
        p3.alignment = PP_ALIGN.CENTER

        # 2. Síntese da Avaliação & Decisão
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "SÍNTESE DA AVALIAÇÃO DA BANCA", "Resumo Executivo e Decisão Editorial")
        add_accent_line(s)
        add_score_badge(s)
        resumo_bullets = _split_bullets(parecer_raw, max_items=5)
        add_bullets(s, resumo_bullets)

        # 3. Estrutura do Documento (Módulo 00)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "ESTRUTURA DO DOCUMENTO", "Conformidade IMRaD, Fluxo Narrativo e Seções (Módulo 00)")
        add_accent_line(s)
        add_bullets(s, _split_bullets(estrutura, max_items=5))

        # 4. Auditoria Metodológica (Módulo 01)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "AUDITORIA METODOLÓGICA", "Desenho Experimental, Validade Interna/Externa e Amostragem (Módulo 01)")
        add_accent_line(s)
        add_bullets(s, _split_bullets(metodologia, max_items=5))

        # 5. Consistência Estatística e Dados (Módulos 01 + 07)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "ANÁLISE ESTATÍSTICA E DADOS", "Consistência Numérica, Testes e Discrepâncias (Módulos 01 & 07)")
        add_accent_line(s)
        dados_text = quantitativa if len(quantitativa) > 100 else metodologia
        add_bullets(s, _split_bullets(dados_text, max_items=5))

        # 6. Checklist Editorial & Ética (Módulo 02)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "CHECKLIST EDITORIAL & ÉTICA", "Aprovação Ética, Conflito de Interesses e Integridade (Módulo 02)")
        add_accent_line(s)
        add_bullets(s, _split_bullets(editorial, max_items=5))

        # 7. Referencial Teórico & SOTA (Módulo 03)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "REFERENCIAL TEÓRICO & SOTA", "Atualidade das Fontes, Estado da Arte e Lacuna (Módulo 03)")
        add_accent_line(s)
        add_bullets(s, _split_bullets(sota, max_items=5))

        # 8. Gaps Lógicos e Argumentativos (Módulo 04)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "GAPS LÓGICOS E ARGUMENTATIVOS", "Falhas de Raciocínio, Não-Sequiturs e Contradições (Módulo 04)", RED)
        add_accent_line(s, color=RED)
        add_bullets(s, _split_bullets(gaps, max_items=5), color=RED)

        # 9. Qualidade da Escrita (Módulo 05)
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "QUALIDADE DA ESCRITA & REGISTRO", "Auditoria Textual, Clareza e Resumo Estatístico de Erros (Módulo 05)", ORANGE)
        add_accent_line(s, color=ORANGE)
        add_bullets(s, _split_bullets(erros_text, max_items=5))

        # 10. Pontos Fortes Identificados
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "PONTOS FORTES IDENTIFICADOS", "Méritos Científicos e Contribuições Validadas pela Banca", GREEN)
        add_accent_line(s, color=GREEN)
        fortes = list(structured.get("pontos_fortes", [])) or _split_bullets(parecer_raw, 4)
        add_bullets(s, fortes, color=GREEN)

        # 11. Fragilidades Principais Detectadas
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "FRAGILIDADES PRINCIPAIS DETECTADAS", "Pontos Críticos que Demandam Correção e Atenção dos Autores", RED)
        add_accent_line(s, color=RED)
        fragilidades = list(structured.get("fragilidades", [])) or _split_bullets(gaps, 4)
        add_bullets(s, fragilidades, color=RED)

        # 12. Recomendações Priorizadas
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "RECOMENDAÇÕES PRIORIZADAS", "Ações Mandatórias e Melhorias Sugeridas para Publicação", TEAL)
        add_accent_line(s, color=TEAL)
        recs = list(structured.get("recomendacoes_obrigatorias", [])) or _split_bullets(parecer_raw, 4)
        add_bullets(s, recs)

        # 13. Auditoria Bibliográfica (Crossref / DOIs se houver)
        if len(bibliografia) > 100:
            s = prs.slides.add_slide(blank_layout)
            add_bg(s, LIGHT_BG)
            add_title(s, "AUDITORIA BIBLIOGRÁFICA & DOIs", "Verificação Automática de Referências e Retratações (Crossref)")
            add_accent_line(s)
            add_bullets(s, _split_bullets(bibliografia, max_items=5))

        # 14. Veredito e Parecer Final
        s = prs.slides.add_slide(blank_layout)
        add_bg(s, NAVY)
        txBox = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.0))
        p = txBox.text_frame.paragraphs[0]
        p.text = "VEREDITO & PARECER FINAL"
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(1.4))
        tf2 = txBox2.text_frame
        p_dec = tf2.paragraphs[0]
        p_dec.text = f"DECISÃO EDITORIAL: {decisao.upper()}"
        p_dec.font.size = Pt(20)
        p_dec.font.bold = True
        p_dec.font.color.rgb = ORANGE if "revis" in decisao.lower() else (GREEN if "accept" in decisao.lower() else RED)
        p_dec.alignment = PP_ALIGN.CENTER

        p_nota = tf2.add_paragraph()
        p_nota.text = f"NOTA FINAL DA BANCA: {nota_str}"
        p_nota.font.size = Pt(18)
        p_nota.font.bold = True
        p_nota.font.color.rgb = WHITE
        p_nota.alignment = PP_ALIGN.CENTER
        p_nota.space_before = Pt(8)

        txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.5))
        p3 = txBox3.text_frame.paragraphs[0]
        p3.text = "Hermes Agent + NotebookLM v6.0 · Análise gerada por IA sujeita a revisão humana."
        p3.font.size = Pt(11)
        p3.font.color.rgb = MUTED
        p3.alignment = PP_ALIGN.CENTER

        out_file = output_dir / "apresentacao_completa.pptx"

    else:  # auditoria
        prs.core_properties.title = "Relatório de Auditoria Peer-Review"
        prs.core_properties.author = "Hermes Agent + NotebookLM v6.0"

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, RED)
        txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "RELATÓRIO DE AUDITORIA CRÍTICA"
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        txBox2 = s.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.8))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = pdf_name
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
        p2.alignment = PP_ALIGN.CENTER
        
        txBox3 = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
        p3 = txBox3.text_frame.paragraphs[0]
        p3.text = f"Auditoria Especializada · Nota: {nota_str} · Decisão: {decisao}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
        p3.alignment = PP_ALIGN.CENTER

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "GAPS LÓGICOS & METODOLÓGICOS", "Problemas Críticos Identificados", RED)
        add_accent_line(s, color=RED)
        add_bullets(s, _split_bullets(gaps, 5))

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "ERROS DE ESCRITA & REGISTRO", "Problemas Textuais e Gramaticais", ORANGE)
        add_accent_line(s, color=ORANGE)
        add_bullets(s, _split_bullets(erros_text, 5))

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, LIGHT_BG)
        add_title(s, "CHECKLIST EDITORIAL", "Conformidade e Diretrizes da Revista")
        add_accent_line(s)
        add_bullets(s, _split_bullets(editorial, 5))

        s = prs.slides.add_slide(blank_layout)
        add_bg(s, NAVY)
        txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "RECOMENDAÇÕES DA BANCA"
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        add_bullets(s, _split_bullets(parecer_raw, 5), y=2.2, h=3, color=RGBColor(0xE8, 0xEA, 0xF6))

        out_file = output_dir / "apresentacao_auditoria.pptx"

    try:
        out_file.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_file))
        size = out_file.stat().st_size
        if size > 1000:
            log(f"  ✓ PPTX gerado: {size} bytes ({out_file.name})")
            return True
        else:
            log(f"  ✗ PPTX muito pequeno: {size} bytes")
            return False
    except Exception as e:
        log(f"  ✗ Erro ao salvar PPTX: {e}")
        return False


# ===== HTML Animado (Mira) =====

def generate_mira_artifact(output_dir: Path, pdf_name: str) -> bool:
    """Gera apresentação animada HTML via padrões mira-animator (D3.js + Tailwind)."""

    parecer_file = output_dir / "06_sintese_parecer.md"
    metodologia_file = output_dir / "01_metodologia.md"
    gaps_file = output_dir / "04_gaps_logicos.md"
    estrutura_file = output_dir / "00_estrutura_documento.md"
    auditoria_file = output_dir / "02_auditoria_editorial.md"
    sota_file = output_dir / "03_sota_referencias.md"
    escrita_file = output_dir / "05_analise_escrita.md"

    parecer = parecer_file.read_text(encoding="utf-8") if parecer_file.exists() else ""
    metodologia = metodologia_file.read_text(encoding="utf-8") if metodologia_file.exists() else ""
    gaps = gaps_file.read_text(encoding="utf-8") if gaps_file.exists() else ""
    estrutura = estrutura_file.read_text(encoding="utf-8") if estrutura_file.exists() else ""
    auditoria = auditoria_file.read_text(encoding="utf-8") if auditoria_file.exists() else ""
    escrita = escrita_file.read_text(encoding="utf-8") if escrita_file.exists() else ""

    from utils import extract_structured_score
    structured = extract_structured_score(parecer)

    nota = str(structured["nota"]) if structured.get("nota") is not None else "N/A"
    decisao = structured.get("decisao") or "N/A"
    _nota_found = nota != "N/A"
    _dec_found = decisao != "N/A"

    for line in parecer.split("\n"):
        # NOTA: regex permissivo — captura qualquer variação de formatação
        if not _nota_found:
            m_nota = re.search(r'(?i)\bNOTA\s*(?:FINAL)?\b[^0-9]*(\d+[.,]?\d*)', line)
            if m_nota and m_nota.group(1):
                nota = m_nota.group(1)
                _nota_found = True
        # DECISÃO: regex permissivo (Accept, Minor/Major Revisions, Reject, etc.)
        if not _dec_found:
            m_dec = re.search(r'(?i)\bDECIS[ÃA]O\s*(?:EDITORIAL|FINAL|RECOMENDADA)?\s*[:=\-—]?\s*(?:\*\*\s*)?(.+?)(?:\s*\*\*)?$', line)
            if m_dec and m_dec.group(1).strip():
                decisao = re.sub(r'\*\*', '', m_dec.group(1)).strip()
                _dec_found = True

    gaps_list = []
    for line in gaps.split("\n"):
        match = re.match(r"^(?:- |\d+\.\s+)(.+)", line.strip())
        if match:
            gap_text = match.group(1).replace("**", "")
            if ":" in gap_text:
                gap_text = gap_text.split(":", 1)[0]
            gaps_list.append(gap_text.strip()[:80])
    gaps_list = gaps_list[:5]

    fortes_list = list(structured.get("pontos_fortes", []))[:4]
    if not fortes_list:
        in_fortes = False
        for line in parecer.split("\n"):
            if re.match(r'(?i)^#{1,3}\s*(?:\d+\.?\d*\s*)?(?:\*\*)?(?:pontos?\s+fortes?|for[çc]as?\s+(?:do|identificad))', line):
                in_fortes = True
                continue
            if in_fortes and re.match(r'^#{1,3}\s', line):
                in_fortes = False

            if in_fortes:
                match = re.match(r"^(?:- |\d+\.\s+)(.+)", line.strip())
                if match:
                    forte_text = match.group(1).replace("**", "")
                    if ":" in forte_text:
                        forte_text = forte_text.split(":", 1)[0]
                    fortes_list.append(forte_text.strip()[:80])
        fortes_list = fortes_list[:4]

    # --- Extract Fragilidades from parecer ---
    fragilidades_list = list(structured.get("fragilidades", []))[:4]
    if not fragilidades_list:
        in_fragilidades = False
        for line in parecer.split("\n"):
            if re.match(r'(?i)^#{1,3}\s*(?:\d+\.?\d*\s*)?(?:\*\*)?(?:fragilidades?|pontos?\s+fracos?|limita[çc][õo]es?|fraquezas?|weakness)', line):
                in_fragilidades = True
                continue
            if in_fragilidades and re.match(r'^#{1,3}\s', line):
                in_fragilidades = False
            if in_fragilidades:
                match = re.match(r"^(?:- |\d+\.\s+)(.+)", line.strip())
                if match:
                    texto = match.group(1).replace("**", "")
                    if ":" in texto:
                        texto = texto.split(":", 1)[0]
                    fragilidades_list.append(texto.strip()[:80])
        fragilidades_list = fragilidades_list[:4]

    # --- Extract Resumo Executivo from parecer ---
    resumo_executivo = ""
    in_resumo = False
    resumo_lines = []
    for line in parecer.split("\n"):
        if re.search(r'(?i)resumo\s*executivo', line):
            in_resumo = True
            continue
        if in_resumo and re.match(r'^#{1,3}\s', line):
            in_resumo = False
        if in_resumo and line.strip():
            resumo_lines.append(line.strip())
    resumo_executivo = " ".join(resumo_lines)[:500]

    # --- Extract methodology steps dynamically from 01_metodologia.md ---
    metodo_steps = []
    for line in metodologia.split("\n"):
        m = re.match(r'^- \*\*(.+?)\*\*', line.strip())
        if m:
            metodo_steps.append(m.group(1).strip()[:50])
    if not metodo_steps:
        for line in metodologia.split("\n"):
            m = re.match(r'^(?:\d+[.)]\s+)(.+)', line.strip())
            if m:
                metodo_steps.append(m.group(1).replace("**", "").strip()[:50])
    if not metodo_steps:
        for line in metodologia.split("\n"):
            m = re.match(r'^#{2,4}\s+(.+)', line.strip())
            if m:
                metodo_steps.append(m.group(1).replace("**", "").strip()[:50])
    if not metodo_steps:
        for line in metodologia.split("\n"):
            m = re.match(r'^(?:[-•*>])\s+(.+)', line.strip())
            if m:
                text = m.group(1).replace("**", "").strip()[:50]
                if len(text) > 5:
                    metodo_steps.append(text)
    metodo_steps = metodo_steps[:5] or ["Design", "Coleta", "Análise", "Conclusão"]

    # --- Extract writing error count from 05_analise_escrita.md ---
    escrita_total_erros = 0
    if escrita:
        # Count table rows (lines starting with |)
        table_rows = [l for l in escrita.split("\n") if l.strip().startswith("|") and "Seção" not in l and "---" not in l]
        escrita_total_erros = len(table_rows)
        # Also try to find explicit count
        for line in escrita.split("\n"):
            m = re.search(r'(\d+)\s+erros?', line, re.IGNORECASE)
            if m:
                escrita_total_erros = int(m.group(1))
                break

    # --- Extract audit checklist items from 02_auditoria_editorial.md ---
    auditoria_items = []
    for line in auditoria.split("\n"):
        # Formato checkbox: - [x] **Item** ou - [ ] **Item**
        m = re.match(r'^- \[([ xX✓✗])\]\s*\*\*(.+?)\*\*', line.strip())
        if m:
            status = "ok" if m.group(1).lower() in ('x', '✓') else "fail"
            auditoria_items.append({"item": m.group(2).strip()[:40], "status": status})
            continue
        # Formato emoji: - ✅ **Item** ou - ❌ **Item**
        m2 = re.match(r'^[-•]\s*([✅✓☑❌✗☒⚠️])\s*\*?\*?(.+?)(?:\*\*)?$', line.strip())
        if m2:
            status = "ok" if m2.group(1) in ('✅', '✓', '☑') else "fail"
            auditoria_items.append({"item": m2.group(2).replace("**", "").strip()[:40], "status": status})
            continue
        # Formato tabela: | Item | Sim/Não |
        m3 = re.match(r'^\|(.+?)\|(.+?)\|', line.strip())
        if m3 and '---' not in m3.group(1):
            item_text = m3.group(1).replace("**", "").strip()
            status_text = m3.group(2).strip().lower()
            if item_text and len(item_text) > 2:
                status = "ok" if any(s in status_text for s in ('sim', 'yes', '✅', '✓', 'ok', 'presente')) else "fail"
                auditoria_items.append({"item": item_text[:40], "status": status})
    auditoria_items = auditoria_items[:6]

    # Diagnóstico: log dos dados extraídos para a apresentação
    log(f"  📊 Dados extraídos para apresentação HTML:")
    log(f"    NOTA: {nota} | DECISÃO: {decisao}")
    log(f"    Gaps: {len(gaps_list)} | Fortes: {len(fortes_list)} | Fragilidades: {len(fragilidades_list)}")
    log(f"    Metodologia steps: {len(metodo_steps)} | Auditoria items: {len(auditoria_items)}")
    log(f"    Resumo executivo: {len(resumo_executivo)} chars | Erros escrita: {escrita_total_erros}")
    if nota == "N/A":
        log("    ⚠️ NOTA não extraída — verifique formato de 06_sintese_parecer.md")
    if decisao == "N/A":
        log("    ⚠️ DECISÃO não extraída — verifique formato de 06_sintese_parecer.md")

    slug = re.sub(r"[^a-z0-9]", "", pdf_name.lower())[:20] or "paper"
    gaps_json = json.dumps(gaps_list[:5] if gaps_list else ["Sem dados"])
    nota_clean = nota.replace(",", ".")
    decisao_short = decisao[:30]
    decisao_60 = decisao[:60]

    # Carregar template Jinja2
    template_dir = Path(__file__).resolve().parent / "templates"
    env = Environment(loader=FileSystemLoader(str(template_dir)), autoescape=select_autoescape(["html", "xml"]))
    template = env.get_template("mira.html")

    html = template.render(
        slug=slug,
        pdf_name=pdf_name,
        nota=nota,
        nota_clean=nota_clean,
        decisao_short=decisao_short,
        decisao_60=decisao_60,
        gaps_json=gaps_json,
        gaps_list=gaps_list,
        fortes_json=json.dumps(fortes_list if fortes_list else ["Nenhum ponto forte destacado"]),
        fortes_list=fortes_list,
        fragilidades_list=fragilidades_list,
        fragilidades_json=json.dumps(fragilidades_list if fragilidades_list else []),
        metodo_steps_json=json.dumps(metodo_steps),
        escrita_total_erros=escrita_total_erros,
        escrita_text=_strip_nb_preamble(escrita)[:400] if escrita else "Análise de escrita não disponível.",
        auditoria_items=auditoria_items,
        auditoria_json=json.dumps(auditoria_items if auditoria_items else []),
        resumo_executivo=resumo_executivo if resumo_executivo else "Resumo executivo não disponível.",
        metodologia_text=_strip_nb_preamble(metodologia)[:500] if metodologia else "Metodologia não analisada.",
        total_slides=7,
    )

    output_file = output_dir / "apresentacao_animada.html"
    output_file.write_text(html, encoding="utf-8")
    size = output_file.stat().st_size

    if size > 1000:
        log(f"  ✓ Apresentação animada HTML gerada: {size} bytes")
        return True
    else:
        log(f"  ✗ HTML muito pequeno: {size} bytes")
        return False



